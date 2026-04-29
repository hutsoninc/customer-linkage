"""
Parse binary input files (PDF, PPTX, DOCX) and print their text content.

Dependencies:
    pip install pymupdf python-pptx python-docx
"""

import sys
import io
from pathlib import Path

# Force UTF-8 output on Windows
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

INPUT_DIR = Path(__file__).parent.parent / "source-materials"

FILES = {
    "pdf": INPUT_DIR / "AHW 2023 Presentation.pdf",
    "pptx": INPUT_DIR / "Customer Data Linkage EXPO Final.pptx",
    "docx": INPUT_DIR / "Customer Data Management Job Aid.docx",
}


def parse_pdf(path: Path) -> str:
    import fitz  # pymupdf
    doc = fitz.open(path)
    pages = []
    for i, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            pages.append(f"--- Page {i} ---\n{text.strip()}")
    return "\n\n".join(pages)


def parse_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = "\n".join(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text.strip():
                    parts.append(text)
        # Speaker notes
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = "\n".join(p.text for p in notes_tf.paragraphs if p.text.strip())
            if notes.strip():
                parts.append(f"[Notes]\n{notes.strip()}")
        if parts:
            slides.append(f"--- Slide {i} ---\n" + "\n\n".join(parts))
    return "\n\n".join(slides)


def parse_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


PARSERS = {
    "pdf": parse_pdf,
    "pptx": parse_pptx,
    "docx": parse_docx,
}


def main():
    targets = sys.argv[1:] if len(sys.argv) > 1 else list(PARSERS.keys())

    for key in targets:
        if key not in FILES:
            print(f"Unknown target: {key}. Choose from: {', '.join(FILES)}")
            continue
        path = FILES[key]
        print(f"\n{'='*70}")
        print(f"FILE: {path.name}")
        print(f"{'='*70}\n")
        try:
            text = PARSERS[key](path)
            print(text if text.strip() else "(no text extracted)")
        except ImportError as e:
            print(f"Missing dependency: {e}")
            print("Install with: pip install pymupdf python-pptx python-docx")
        except Exception as e:
            print(f"Error: {e}")


if __name__ == "__main__":
    main()
